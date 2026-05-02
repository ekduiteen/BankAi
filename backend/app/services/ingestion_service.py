import os
import uuid
from pypdf import PdfReader
from docx import Document as DocxDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client.models import PointStruct
from sqlmodel import Session

from ..models.document import Document, DocumentChunk
from .embedding_service import generate_embeddings
from .qdrant_service import upload_points


def extract_text(file_path: str, file_type: str) -> str:
    text = ""
    ft = file_type.lower()
    if ft == 'pdf':
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif ft == 'docx':
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif ft == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    elif ft in ['xlsx', 'xls']:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            text += f"\n--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                vals = [str(c) if c is not None else "" for c in row]
                text += " | ".join(vals) + "\n"
        wb.close()
    elif ft in ['pptx', 'ppt']:
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ft in ['jpg', 'jpeg', 'png']:
        import base64
        from .llm_service import call_vision_llm
        with open(file_path, 'rb') as f:
            image_b64 = base64.b64encode(f.read()).decode('utf-8')
        description = call_vision_llm(
            "Please provide a detailed text description of this document or image. "
            "Extract any text you see. If it's a financial document, list key numbers and dates.",
            image_b64
        )
        text = f"Image/Photo Description:\n{description}"
    return text


def auto_catalog(text: str, filename: str) -> dict:
    """
    Use LLM to automatically classify document type, department, and suggest title.
    Returns dict with keys: document_type, department, title.
    Falls back gracefully if LLM can't parse response.
    """
    import json
    from .llm_service import call_llm

    try:
        prompt = f"""Analyze this document excerpt and classify it. Respond ONLY as valid JSON with no other text.

Extract and return:
1. "document_type": exactly one of [policy, procedure, manual, compliance, report, circular, directive, other]
2. "department": one of [Credit, Treasury, Operations, HR, Compliance, Audit, IT, Risk, General]
3. "title": a concise document title (max 10 words)

Document filename: {filename}

First 3000 characters of content:
{text[:3000]}

Respond ONLY as JSON, no markdown, no explanation."""

        response = call_llm(prompt).strip()

        # Try to extract JSON from response (handles models that may add markdown formatting)
        if "```json" in response:
            response = response.split("```json")[1].split("```")[0].strip()
        elif "```" in response:
            response = response.split("```")[1].split("```")[0].strip()

        data = json.loads(response)

        return {
            "document_type": data.get("document_type", "other"),
            "department": data.get("department", "General"),
            "title": data.get("title", filename)
        }
    except Exception as e:
        # Gracefully fall back to defaults if LLM fails
        return {
            "document_type": "other",
            "department": "General",
            "title": filename
        }


def process_document(document_id: int):
    """Background task — creates its own DB session so it is safe to run after the request ends."""
    from ..db.session import engine

    with Session(engine) as db:
        try:
            doc = db.get(Document, document_id)
            if not doc:
                return

            doc.status = "extracting_text"
            doc.processing_progress = 10
            doc.processing_message = "Reading document text..."
            db.add(doc)
            db.commit()

            # 1. Extract text
            text = extract_text(doc.file_path, doc.file_type)

            # Auto-catalog the document
            doc.status = "cataloging"
            doc.processing_progress = 25
            doc.processing_message = "Classifying document..."
            db.add(doc)
            db.commit()

            catalog = auto_catalog(text, doc.file_name)
            # Only update if current type is "other" or fields are empty
            if doc.document_type == "other":
                doc.document_type = catalog["document_type"]
            if not doc.department:
                doc.department = catalog["department"]
            if doc.title == doc.file_name:  # title was auto-filled with filename
                doc.title = catalog["title"]
            db.add(doc)
            db.commit()

            doc.status = "chunking"
            doc.processing_progress = 40
            doc.processing_message = "Preparing document for search..."
            db.add(doc)
            db.commit()

            # 2. Split into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )
            chunks = text_splitter.split_text(text)

            if not chunks:
                doc.status = "failed"
                doc.processing_message = "Failed to extract text."
                db.add(doc)
                db.commit()
                return

            doc.status = "embedding"
            doc.processing_progress = 60
            doc.processing_message = "Creating secure document index..."
            db.add(doc)
            db.commit()

            # 3. Generate embeddings in batches
            BATCH = 64
            embeddings = []
            for i in range(0, len(chunks), BATCH):
                embeddings.extend(generate_embeddings(chunks[i:i + BATCH]))

            # 4. Store in Qdrant and DB (bulk insert)
            points = []
            chunk_records = []
            for i, (chunk_text, embedding) in enumerate(zip(chunks, embeddings)):
                point_id = str(uuid.uuid4())
                chunk_records.append(DocumentChunk(
                    bank_id=doc.bank_id,
                    document_id=doc.id,
                    chunk_index=i,
                    chunk_text=chunk_text,
                    qdrant_point_id=point_id
                ))
                points.append(PointStruct(
                    id=point_id,
                    vector=embedding,
                    payload={
                        "bank_id": doc.bank_id,
                        "document_id": doc.id,
                        "chunk_index": i,
                        "text": chunk_text,
                        "document_scope": doc.document_scope,
                        "session_id": doc.session_id,
                    }
                ))

            db.add_all(chunk_records)
            db.commit()

            doc.status = "indexing"
            doc.processing_progress = 90
            doc.processing_message = "Adding document to session..."
            db.add(doc)
            db.commit()

            # Upload to Qdrant in batches
            QDRANT_BATCH = 100
            for i in range(0, len(points), QDRANT_BATCH):
                upload_points(points[i:i + QDRANT_BATCH])

            # 5. Generate summary
            try:
                from .llm_service import call_llm
                summary_prompt = f"Summarize the following document concisely in 2-3 sentences:\n\n{text[:4000]}"
                doc.summary = call_llm(summary_prompt)
            except Exception:
                doc.summary = "Summary generation failed."

            doc.status = "ready"
            doc.processing_progress = 100
            doc.processing_message = "Document is ready. You can ask follow-up questions."
            db.add(doc)
            db.commit()

        except Exception as e:
            try:
                doc = db.get(Document, document_id)
                if doc:
                    doc.status = "failed"
                    doc.processing_message = f"Error: {str(e)}"
                    db.add(doc)
                    db.commit()
            except Exception:
                pass
